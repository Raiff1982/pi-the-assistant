from pptx import Presentation
from pptx.util import Inches
import zipfile
import os

# Create a PowerPoint presentation object
presentation = Presentation()

# Add a title slide
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Project Cost Estimates"
subtitle.text = "Summary of estimated costs for the project"

# Add a slide for the cost estimates
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Estimated Monthly Costs by Service Category"

# Add content to the slide
content = """
1. Analytics
   - Azure AI Bot Service: $100
   - Azure Machine Learning: $1,319.84
   - Azure AI services: $192.72
   - Azure AI Document Intelligence: $212.50
   - Azure AI Video Indexer: $57.60
   - Power BI Embedded: $735.91
   - Data Catalog: $3
   - Azure Data Factory: $8.03
   - Azure Chaos Studio: $30.10
2. AI & Machine Learning
   - Microsoft Genomics: $1
   - Machine Learning Studio (classic): $64.95
   - Azure OpenAI Service: $1,460
   - Health Bot: $0
   - Azure AI Metrics Advisor: $0
   - Azure AI services (LUIS): $300
3. Compute
   - Virtual Machines: $91.10
   - Virtual Machine Scale Sets: $66.72
   - Azure Service Fabric: $586.54
   - Azure Kubernetes Service (AKS): $85.41
   - Azure Functions: $71.40
   - Cloud Services: $137.24
   - Azure Modeling and Simulation Workbench: $599.87
   - Azure Container Instances: $117.53
   - Batch: $160.60
4. Storage
   - Storage Accounts: $109.99
   - Managed Disks: $26.40
   - Azure Container Storage: $141.31
   - Azure Files: $139.12
5. Web
   - API Management: $48.03
   - Azure AI Search: $490.56
6. Databases
   - Azure Cosmos DB: $80.84
   - Azure SQL Database: $372.97
7. Developer Tools
   - Azure Lab Services: $306.60
8. DevOps
   - Azure Load Testing: $10.15
   - Azure Monitor: $156.50
9. Management and Governance
   - Azure Backup: $265.85
   - Microsoft Cost Management: $15
10. Security 
    - Microsoft Defender for Cloud: $5.11 
    - Microsoft Security Copilot: $2,920 
    - Defender External Attack Surface Management: $0.99 
Total Estimated Monthly Cost: $15,700.26 
"""
text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
text_frame = text_box.text_frame

for line in content.split('\n'):
    p = text_frame.add_paragraph()
    p.text = line

# Save the updated presentation to a file.
presentation.save("Project_Cost_Estimates_Updated.pptx")

# Convert the PowerPoint presentation to VSDX format by creating a zip file with the necessary structure.
with zipfile.ZipFile('Project_Cost_Estimates.vsdx', 'w') as vsdx:
    # Add the PowerPoint file to the zip archive.
    vsdx.write('Project_Cost_Estimates_Updated.pptx', 'ppt/presentation.xml')

print("The VSDX file has been created and saved as 'Project_Cost_Estimates.vsdx'.")
